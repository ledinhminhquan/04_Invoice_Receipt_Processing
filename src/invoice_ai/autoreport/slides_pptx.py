"""Auto-generate the ~12-slide PPTX submission deck from artifacts."""

from __future__ import annotations

from pathlib import Path
from typing import List, Optional, Tuple

from ..config import AppConfig, artifacts_dir
from ..logging_utils import get_logger, utc_stamp
from .artifact_loader import load_all_artifacts
from .charts import build_all_charts

logger = get_logger(__name__)

# (title, [bullets], chart_key|None)
_SLIDES: List[Tuple[str, List[str], Optional[str]]] = [
    ("Invoice & Receipt Processing System",
     ["Document-AI KIE: image/PDF -> validated structured JSON.",
      "Offline-first, agentic, with arithmetic reconciliation + human-in-the-loop.",
      "NLP in Industry — Final Assignment, Project #4."], None),
    ("Business Problem & Motivation",
     ["Manual invoice data entry is slow, costly and error-prone.",
      "OCR-only tools don't validate; pure-LLM tools silently hallucinate totals.",
      "A wrong total flows straight into the ledger — undetected."], None),
    ("Proposed Solution",
     ["Local LayoutLMv3/Donut extractor (no mandatory cloud API).",
      "Validation engine reconciles subtotal + tax == total (Decimal money).",
      "That GPT-4o-vision reference call is demoted to one optional fallback tool."], None),
    ("System Architecture",
     ["classify -> OCR -> extract fields -> line items -> validate -> normalize -> decide.",
      "Single AgentState blackboard threaded through 7 tools, fully traced.",
      "Born-digital PDF (pdfplumber) vs scanned (pdf2image + OCR) router."], None),
    ("Data Overview",
     ["mp-02/sroie (receipt KIE), naver-clova-ix/cord-v2 (line items, CC-BY-4.0).",
      "nielsr/funsd-layoutlmv3 (forms), katanaml invoices (MIT).",
      "No large data committed; bbox normalised to 0-1000 for LayoutLMv3."], None),
    ("Model & Evaluation Results",
     ["LayoutLMv3-base (accuracy) / LiLT (MIT, commercial) / Donut (line items).",
      "Baselines: regex/heuristics + bert+bbox. Metric: entity-level seqeval F1.",
      "Trained model must beat the regex floor to justify its cost."], "metrics"),
    ("Agentic AI Component",
     ["Three decision points: doc-type/quality routing, validation, final confidence gate.",
      "Worked example: printed total 1,860 vs subtotal+tax 2,160 -> flagged (delta -300).",
      "Auto-approve only when reconciled AND confident; else human review."], None),
    ("Deployment Overview",
     ["FastAPI /extract /classify /batch /review-queue /metrics + Gradio highlight demo.",
      "Docker (tesseract + poppler) on HF Space, port 7860; model_version echoed.",
      "Latency /extract ~250-500 ms GPU / ~0.8-1.5 s CPU per page."], "latency"),
    ("Ethics, Privacy & Risks",
     ["Offline-first = invoices never leave the org (privacy by design).",
      "PII redaction, Decimal money, human-in-the-loop, bbox provenance for audit.",
      "Reconciliation guards against silent financial errors and tampering."], None),
    ("Continual Learning & Monitoring",
     ["Human review-queue corrections become retraining labels.",
      "Monitor field-F1, needs-review rate, OCR-confidence drift; canary by review-rate.",
      "/metrics exposes latency, confidence quantiles, needs-review counts."], None),
    ("Key Takeaways & Future Work",
     ["Strictly dominates the GPT-4o-only reference: offline, validated, multi-page, HITL.",
      "Trained LayoutLMv3 grounds every field to a box for auditable review.",
      "Future: LiLT multilingual, Donut line-items, active learning from the queue."], None),
    ("Q&A",
     ["Thank you — Le Dinh Minh Quan (23127460), package: invoice_ai.",
      "Demo: HF Space (Docker, 7860) + FastAPI /extract.",
      "Recap: extract -> validate -> decide, offline + agentic."], None),
]


def generate_slides(cfg: AppConfig, title: str = None, author: str = None, out_path=None) -> Path:
    title = title or cfg.project_title
    author = author or cfg.author
    stamp = utc_stamp()
    out = Path(out_path) if out_path else artifacts_dir() / "submission" / f"submission-{stamp}" / "slides.pptx"
    out.parent.mkdir(parents=True, exist_ok=True)
    charts = build_all_charts(load_all_artifacts(), out_dir=str(out.parent / "_figures"))

    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
    except Exception as exc:
        logger.warning("python-pptx unavailable (%s); markdown fallback.", exc)
        return _md_fallback(out, title, author)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    NAVY = RGBColor(0x1F, 0x2A, 0x44)
    ACCENT = RGBColor(0x2E, 0x6F, 0xF2)

    for i, (heading, bullets, chart_key) in enumerate(_SLIDES):
        slide = prs.slides.add_slide(blank)
        htb = slide.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(12), Inches(1))
        hp = htb.text_frame.paragraphs[0]
        hp.text = heading if i else title
        hp.font.size = Pt(30 if i else 38)
        hp.font.bold = True
        hp.font.color.rgb = NAVY
        body_w = Inches(6.6) if chart_key and chart_key in charts else Inches(12)
        body = slide.shapes.add_textbox(Inches(0.6), Inches(1.6), body_w, Inches(5))
        tf = body.text_frame
        tf.word_wrap = True
        for j, b in enumerate(bullets):
            para = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
            para.text = "•  " + b
            para.font.size = Pt(16)
            para.space_after = Pt(10)
        if chart_key and chart_key in charts:
            try:
                slide.shapes.add_picture(str(charts[chart_key]), Inches(7.4), Inches(1.7), width=Inches(5.3))
            except Exception:
                pass
        sub = slide.shapes.add_textbox(Inches(0.5), Inches(7.0), Inches(12), Inches(0.4))
        sp = sub.text_frame.paragraphs[0]
        sp.text = f"Invoice & Receipt Processing · {author} · slide {i+1}/{len(_SLIDES)}"
        sp.font.size = Pt(9)
        sp.font.color.rgb = ACCENT

    prs.save(str(out))
    logger.info("Wrote slide deck (%d slides) -> %s", len(_SLIDES), out)
    return out


def _md_fallback(out: Path, title, author) -> Path:
    md = out.with_suffix(".md")
    lines = [f"# {title}", f"_{author} — Project #4_", ""]
    for i, (h, bullets, _) in enumerate(_SLIDES):
        lines.append(f"## Slide {i+1} — {h}")
        lines += [f"- {b}" for b in bullets] + [""]
    md.write_text("\n".join(lines), encoding="utf-8")
    return md


__all__ = ["generate_slides"]
